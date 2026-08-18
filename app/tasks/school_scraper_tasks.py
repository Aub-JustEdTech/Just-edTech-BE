"""Celery tasks for ingesting scraped media into the document pipeline.

Scrape orchestration (discovery, run/job tracking) is handled offline via
scripts; this module covers materialize (download or transcribe) -> Document
-> vectors, plus the batch sweep that walks confirmed school URLs.

Cost model, cheapest gate first:

* documents (PDF/DOCX/PPTX) — downloaded and text-extracted locally, free;
* YouTube with captions — captions fetched, free;
* everything else — AssemblyAI at ~$0.23/audio-hour, capped by duration.
"""

from __future__ import annotations

import logging
import shutil
import tempfile
from dataclasses import dataclass, field
from datetime import datetime, timezone
from pathlib import Path

from sqlalchemy.ext.asyncio import AsyncSession

from app.celery_app import celery_app
from app.core.config import settings
from app.db.connector import AsyncSessionLocal
from app.models.school import ScrapedMedia
from app.services.transcription.exceptions import (
    TerminalTranscriptionError,
)
from app.services.transcription.schemas import TranscriptResult
from app.tasks.loop_utils import get_event_loop

logger = logging.getLogger(__name__)

# Media types that go down the transcription path rather than text extraction.
AV_MEDIA_TYPES = ("audio", "video", "youtube")


@dataclass(slots=True)
class MediaPayload:
    """What materializing one ScrapedMedia produced."""

    text: str = ""
    transcript: TranscriptResult | None = None
    content_hash: str | None = None
    size_bytes: int | None = None
    duration_seconds: int | None = None
    raw_bytes: bytes | None = None
    extra_metadata: dict = field(default_factory=dict)


@celery_app.task(
    name="app.tasks.school_scraper_tasks.ingest_scraped_media",
    bind=True,
    max_retries=3,
)
def ingest_scraped_media(self, scraped_media_id: int):
    """Download/transcribe a scraped media item and ingest into the pipeline."""
    try:
        loop = get_event_loop()
        return loop.run_until_complete(_ingest_scraped_media_async(scraped_media_id))
    except TerminalTranscriptionError:
        # Already recorded as a status by the async impl; retrying is pointless.
        raise
    except Exception as exc:
        logger.error(
            "ingest_scraped_media failed for %s: %s", scraped_media_id, str(exc)
        )
        raise self.retry(
            exc=exc, countdown=60 * (2**self.request.retries)
        ) from exc


async def _ingest_scraped_media_async(scraped_media_id: int) -> dict:
    from app.crud.schools import (
        get_scraped_media_by_content_hash,
        update_scraped_media,
    )
    from app.services.web_scraper.year_filter import evaluate_media_year_async

    async with AsyncSessionLocal() as db:
        sm = await db.get(ScrapedMedia, scraped_media_id)
        if not sm:
            logger.warning("ScrapedMedia %s not found, skipping", scraped_media_id)
            return {"scraped_media_id": scraped_media_id, "error": "not found"}

        inferred_year, should_process, skip_reason = await evaluate_media_year_async(
            url=sm.source_media_url,
            filename=sm.original_name,
            source_page_url=sm.source_page_url,
        )
        if not should_process:
            await update_scraped_media(
                db,
                sm.id,
                status="skipped_year",
                doc_year=inferred_year,
                error_message=skip_reason,
            )
            return {
                "scraped_media_id": scraped_media_id,
                "status": "skipped_year",
                "doc_year": inferred_year,
            }

        await update_scraped_media(db, sm.id, status="downloading")

        if inferred_year is not None:
            sm.doc_year = inferred_year
            await db.flush()

        # Temp dir is load-bearing: celery-scraper shares the temp_uploads
        # volume with the documents worker, so a leaked multi-GB file takes
        # down both. Always cleaned, even on failure.
        Path(settings.SCHOOL_SCRAPER_MEDIA_TEMP_DIR).mkdir(parents=True, exist_ok=True)
        workdir = Path(
            tempfile.mkdtemp(dir=settings.SCHOOL_SCRAPER_MEDIA_TEMP_DIR)
        )

        try:
            try:
                payload = await _materialize_media(sm, workdir)
            except TerminalTranscriptionError as exc:
                # Deterministic: retrying re-does identical work and fails the
                # same way. Record and return WITHOUT raising.
                logger.warning(
                    "Terminal failure for scraped_media %s (%s): %s",
                    scraped_media_id,
                    exc.status,
                    exc,
                )
                await update_scraped_media(
                    db,
                    sm.id,
                    status=exc.status,
                    error_message=str(exc),
                )
                return {
                    "scraped_media_id": scraped_media_id,
                    "status": exc.status,
                }
            except Exception as exc:
                logger.exception("Ingest failed for scraped_media %s", scraped_media_id)
                await update_scraped_media(
                    db,
                    sm.id,
                    status="failed",
                    error_message=str(exc),
                )
                raise

            # An empty transcript must not create a Document: stage 2 raises
            # on empty text, which would strand the Document in PROCESSING.
            if not payload.text.strip():
                logger.warning(
                    "Empty transcript for scraped_media %s; not creating a Document",
                    scraped_media_id,
                )
                await update_scraped_media(
                    db,
                    sm.id,
                    status="no_transcript",
                    error_message="transcript was empty",
                )
                return {
                    "scraped_media_id": scraped_media_id,
                    "status": "no_transcript",
                }

            content_h = payload.content_hash
            if content_h:
                existing = await get_scraped_media_by_content_hash(
                    db, sm.school_id, content_h
                )
                if existing and existing.id != sm.id:
                    await update_scraped_media(
                        db,
                        sm.id,
                        status="skipped_duplicate",
                        content_hash=content_h,
                    )
                    return {
                        "scraped_media_id": scraped_media_id,
                        "status": "skipped_duplicate",
                    }

            await update_scraped_media(
                db,
                sm.id,
                status="ingesting",
                content_hash=content_h,
            )

            # Everything past this point runs AFTER transcription has already
            # been paid for. An S3 or DB failure here must NOT propagate to the
            # wrapper's self.retry(), or the retry re-transcribes and re-bills
            # the same media up to three more times. Record and stop.
            try:
                document_id = await _create_document_and_enqueue(db, sm, payload)
            except Exception as exc:  # noqa: BLE001
                logger.exception(
                    "Persisting the transcript failed for scraped_media %s AFTER "
                    "transcription was paid for; not retrying to avoid re-billing",
                    scraped_media_id,
                )
                await update_scraped_media(
                    db,
                    sm.id,
                    status="failed",
                    error_message=f"post-transcription persist failed: {exc}",
                    duration_seconds=payload.duration_seconds,
                    size_bytes=payload.size_bytes,
                )
                return {
                    "scraped_media_id": scraped_media_id,
                    "status": "failed",
                    "error": "post-transcription persist failed",
                }

            await update_scraped_media(
                db,
                sm.id,
                status="completed",
                document_id=document_id,
                duration_seconds=payload.duration_seconds,
                size_bytes=payload.size_bytes,
                ingested_at=datetime.now(timezone.utc),
            )
            return {
                "scraped_media_id": scraped_media_id,
                "status": "completed",
                "document_id": document_id,
            }
        finally:
            shutil.rmtree(workdir, ignore_errors=True)


async def _materialize_media(sm, workdir: Path) -> MediaPayload:
    """Produce text (and, for A/V, a transcript) for one ScrapedMedia item."""
    import hashlib

    import httpx

    from app.services.transcription.service import transcription_service
    from app.services.transcription.youtube import extract_youtube_id
    from app.services.web_scraper.board_platforms import (
        fetch_document_via_playwright_session,
        is_board_platform_url,
    )

    # --- YouTube: captions first, always free ---
    if sm.media_type == "youtube":
        if not settings.SCHOOL_SCRAPER_YOUTUBE_TRANSCRIPT_ENABLED:
            from app.services.transcription.exceptions import (
                NoTranscriptAvailableError,
            )

            raise NoTranscriptAvailableError(
                "YouTube transcripts are disabled "
                "(SCHOOL_SCRAPER_YOUTUBE_TRANSCRIPT_ENABLED=False)"
            )

        transcript = await transcription_service.transcribe_youtube(
            sm.source_media_url, workdir=workdir
        )
        video_id = extract_youtube_id(sm.source_media_url) or sm.source_media_url
        return MediaPayload(
            text=transcript.text,
            transcript=transcript,
            # Reuses the existing content_hash column and its unique
            # constraint, so the same video embedded on three pages ingests
            # once — with no schema change.
            content_hash=hashlib.sha256(f"youtube:{video_id}".encode()).hexdigest(),
            duration_seconds=transcript.duration_seconds,
            size_bytes=transcript.source_size_bytes,
        )

    # --- Audio / video: the only path that can cost money ---
    if sm.media_type in ("audio", "video"):
        transcript = await transcription_service.transcribe_media_url(
            sm.source_media_url, workdir=workdir
        )
        return MediaPayload(
            text=transcript.text,
            transcript=transcript,
            duration_seconds=transcript.duration_seconds,
            # Read from the container header during the pre-spend probe, so
            # this is populated even under url_direct where nothing is
            # downloaded.
            size_bytes=transcript.source_size_bytes,
        )

    # --- Documents: board platforms need a Playwright session; others use httpx ---
    # Board-platform downloads require cookies/referrer from a real browser
    # session. A cold httpx.GET from this Celery worker is typically redirected
    # to a login/error page. Non-document media on these platforms is rare;
    # if it occurs, fall through to the httpx/transcription paths above.
    board_doc_url = sm.source_media_url or sm.source_page_url or ""
    if sm.media_type == "document" and is_board_platform_url(board_doc_url):
        raw = await fetch_document_via_playwright_session(
            sm.source_page_url,
            sm.source_media_url,
        )
        return MediaPayload(
            text=_extract_text_from_document(raw, sm.file_extension),
            content_hash=hashlib.sha256(raw).hexdigest(),
            size_bytes=len(raw),
            raw_bytes=raw,
        )

    async with httpx.AsyncClient(
        timeout=settings.WEB_SCRAPER_TIMEOUT_SECONDS,
        headers={"User-Agent": settings.SCHOOL_SCRAPER_USER_AGENT},
        follow_redirects=True,
    ) as client:
        resp = await client.get(sm.source_media_url)
        resp.raise_for_status()
        raw = resp.content

    return MediaPayload(
        text=_extract_text_from_document(raw, sm.file_extension),
        content_hash=hashlib.sha256(raw).hexdigest(),
        size_bytes=len(raw),
        raw_bytes=raw,
    )


def _extract_text_from_document(raw: bytes, ext: str | None) -> str:
    import io
    import os
    import tempfile

    ext = (ext or "").lower().lstrip(".")
    if ext == "pdf":
        import fitz  # pymupdf

        with fitz.open(stream=raw, filetype="pdf") as doc:
            return "\n".join(page.get_text() for page in doc)
    if ext == "docx":
        import docx

        d = docx.Document(io.BytesIO(raw))
        return "\n".join(p.text for p in d.paragraphs)
    if ext == "pptx":
        from pptx import Presentation

        presentation = Presentation(io.BytesIO(raw))
        parts: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        run_text = "".join(run.text for run in paragraph.runs)
                        if run_text or paragraph.text:
                            parts.append(run_text or paragraph.text)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text:
                                parts.append(cell.text)
        return "\n".join(parts)
    if ext in ("doc", "xlsx", "xls", "txt", "text", "md"):
        suffix = f".{ext}"
        with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
            tmp.write(raw)
            tmp_path = tmp.name
        try:
            from app.services.document_processing.factory import ProcessorFactory

            processor = ProcessorFactory.get_processor(tmp_path)
            return processor.extract_text(tmp_path)
        finally:
            os.unlink(tmp_path)
    try:
        return raw.decode("utf-8", errors="ignore")
    except Exception:  # noqa: BLE001
        return ""


def _get_s3_manager():
    from app.utils.s3 import S3Manager

    return S3Manager(
        bucket_name=settings.S3_BUCKET_NAME,
        region_name=settings.S3_REGION,
        aws_access_key_id=settings.AWS_ACCESS_KEY_ID,
        aws_secret_access_key=settings.AWS_SECRET_ACCESS_KEY,
    )


async def _create_document_and_enqueue(
    db: AsyncSession,
    sm,
    payload: MediaPayload,
) -> int | None:
    """Persist the artifacts, create the Document, enqueue the pipeline."""
    from app.models.documents import Document, ProcessingStatus
    from app.models.school import School
    from app.utils.school_calendar import (
        derive_quarter_month,
        derive_school_year,
    )

    s3 = _get_s3_manager()
    key_prefix = (
        f"{settings.SCHOOL_SCRAPER_S3_PREFIX}"
        f"tenants/{sm.tenant_id}/schools/{sm.school_org_code}/"
        f"{sm.media_type}/{payload.content_hash or sm.url_hash}"
    )

    if payload.transcript is not None:
        # One text artifact. Timestamps and speaker labels live in each line's
        # prefix, so both survive into chunking. See TranscriptProcessor.
        document_type = ".transcript"
        s3_key_text = f"{key_prefix}/transcript.txt"
        s3_url_text = await s3.upload_file_object(
            payload.transcript.to_text_document().encode("utf-8"),
            s3_key_text,
        )
    else:
        # Leading dot is required: the pipeline builds "{uuid}{document_type}"
        # and ProcessorFactory keys on Path(...).suffix.
        document_type = ".txt"
        s3_key_text = f"{key_prefix}/transcript.txt"
        s3_url_text = await s3.upload_file_object(
            payload.text.encode("utf-8"), s3_key_text
        )

    # Raw archive exists only for documents, which are downloaded anyway.
    # Under TRANSCRIPTION_AUDIO_MODE=url_direct the worker never holds the
    # audio/video bytes — AssemblyAI fetches the URL itself — so there is
    # nothing to archive. That is the documented trade-off of url_direct;
    # source_media_url remains the pointer to playable media.
    s3_key_raw = None
    if payload.raw_bytes is not None:
        ext = (sm.file_extension or "bin").lstrip(".")
        s3_key_raw = f"{key_prefix}/{sm.original_name or f'file.{ext}'}"
        await s3.upload_file_object(payload.raw_bytes, s3_key_raw)

    # Prefer the raw binary for document types so the pipeline (including OCR)
    # can process the real file instead of our own no-OCR local extraction.
    # Audio/video/YouTube stay on the transcript text artifact.
    _RAW_DOC_EXTS = {"pdf", "docx", "doc", "pptx", "xlsx", "xls"}
    file_ext = (sm.file_extension or "bin").lstrip(".").lower()
    use_raw_document = (
        sm.media_type == "document"
        and s3_key_raw is not None
        and file_ext in _RAW_DOC_EXTS
    )
    if use_raw_document:
        doc_s3_url = f"s3://{settings.S3_BUCKET_NAME}/{s3_key_raw}"
        doc_type = f".{file_ext}"
    else:
        doc_s3_url = s3_url_text
        doc_type = document_type

    # Resolve the school's state (2-letter abbreviation) so it can be
    # denormalized onto the Document row + source_metadata. Falls back to
    # 'MA' if the school row is missing (V1 corpus is MA-only).
    school = await db.get(School, sm.school_id)
    state = (school.state if school else None) or "MA"

    # Derive school_year + quarter_month from the scraped meeting_date if
    # present. DocClassifier (step 2.6) may overwrite meeting_date from the
    # LLM and re-derive these; setting them now means non-LLM paths still
    # have correct values.
    school_year: str | None = None
    quarter_month: str | None = None
    if sm.meeting_date:
        school_year = derive_school_year(sm.meeting_date)
        quarter_month = derive_quarter_month(sm.meeting_date)

    transcript = payload.transcript
    doc = Document(
        name=sm.original_name or sm.source_media_url,
        doc_id=(
            f"school-{sm.school_org_code}-"
            f"{payload.content_hash or sm.url_hash[:16]}"
        ),
        s3_url=doc_s3_url,
        tenant_id=sm.tenant_id,
        document_type=doc_type,
        processing_status=ProcessingStatus.PENDING,
        source_type="school_scraper",
        content_hash=payload.content_hash,
        # Heatmap V1 doc-level denorm (spec: Heatmap Ingest Metadata v1).
        state=state,
        district_name=sm.school_name,
        school_year=school_year,
        quarter_month=quarter_month,
        meeting_date=sm.meeting_date,
        source_metadata={
            "scraped_media_id": sm.id,
            "school_id": sm.school_id,
            "school_org_code": sm.school_org_code,
            "school_name": sm.school_name,
            "district_type": sm.district_type,
            "state": state,
            "source_page_url": sm.source_page_url,
            "source_media_url": sm.source_media_url,
            "media_type": sm.media_type,
            "document_type": sm.document_type,
            "meeting_date": sm.meeting_date.isoformat() if sm.meeting_date else None,
            "school_year": school_year,
            "quarter_month": quarter_month,
            "doc_year": sm.doc_year,
            "scraped_at": sm.scraped_at.isoformat() if sm.scraped_at else None,
            # Provenance needed to resolve a citation back to playable media.
            "transcript_source": transcript.source if transcript else None,
            "speech_model": transcript.speech_model if transcript else None,
            "caption_kind": transcript.caption_kind if transcript else None,
            "duration_seconds": payload.duration_seconds,
            "s3_key_raw": s3_key_raw,
        },
    )
    db.add(doc)
    await db.flush()

    sm.s3_key_raw = s3_key_raw
    sm.s3_key_text = s3_key_text
    if payload.size_bytes is not None:
        sm.size_bytes = payload.size_bytes

    from app.models.processing_jobs import DocumentProcessingJob, JobStatus

    job = DocumentProcessingJob(
        document_id=doc.id,
        status=JobStatus.PENDING,
        processor_type=doc.document_type,
    )
    db.add(job)
    await db.commit()

    try:
        from app.tasks.document_pipeline import process_document_pipeline

        process_document_pipeline.delay(doc.id, job.id)
    except Exception:  # noqa: BLE001
        logger.exception("Failed to enqueue document processing for doc %s", doc.id)

    return doc.id


@celery_app.task(
    name="app.tasks.school_scraper_tasks.sweep_school_media",
    bind=True,
    max_retries=3,
)
def sweep_school_media(self, school_ids: list[int] | None = None) -> dict:
    """Walk every active school_scrape_urls row, persist media, enqueue new rows.

    Does not re-run URL discovery — only the human-confirmed scrapable URLs
    paired with each school are crawled. Deliberately NOT in beat_schedule:
    run it manually and confirm the created/skipped counts look right before
    letting it fire unattended against several hundred district sites.
    """
    try:
        loop = get_event_loop()
        return loop.run_until_complete(_sweep_school_media_async(school_ids))
    except Exception as exc:
        logger.error("sweep_school_media failed: %s", str(exc))
        raise self.retry(
            exc=exc, countdown=60 * (2**self.request.retries)
        ) from exc


async def _sweep_school_media_async(school_ids: list[int] | None) -> dict:
    from sqlalchemy import select

    from app.crud.schools import bulk_create_scraped_media
    from app.models.school import School, SchoolScrapeUrl
    from app.services.web_scraper.school_scraper_service import SchoolScraperService

    totals = {
        "schools": 0,
        "found": 0,
        "created": 0,
        "skipped": 0,
        "enqueued": 0,
        "av_found": 0,
        "documents_found": 0,
        # Counted separately so a run of all-zeros cannot be mistaken for
        # "nothing new was published".
        "scrape_failures": 0,
    }
    failed_urls: list[str] = []

    async with AsyncSessionLocal() as db:
        stmt = select(SchoolScrapeUrl).where(SchoolScrapeUrl.is_active.is_(True))
        if school_ids:
            stmt = stmt.where(SchoolScrapeUrl.school_id.in_(school_ids))
        scrape_urls = (await db.execute(stmt)).scalars().all()

        # `async with` matters twice: it closes the httpx client and any
        # auto-launched Chromium (otherwise leaked per sweep), and its
        # __aenter__ pre-launches the browser so it is reused across all
        # districts instead of relaunched per site.
        async with SchoolScraperService() as service:
            for scrape_url in scrape_urls:
                school = await db.get(School, scrape_url.school_id)
                if not school:
                    continue
                totals["schools"] += 1

                try:
                    result = await service.scrape_media_files(
                        page_url=scrape_url.url,
                        crawl_depth=scrape_url.crawl_depth,
                    )
                except Exception:
                    # One unreachable district must not abort the sweep, but it
                    # must be visible: a silent skip looks identical to a site
                    # with no new meetings.
                    logger.exception("Sweep scrape failed for %s", scrape_url.url)
                    totals["scrape_failures"] += 1
                    failed_urls.append(scrape_url.url)
                    continue

                media_files = result.get("media_files", [])
                totals["found"] += len(media_files)
                for media in media_files:
                    media_type = media.get("media_type")
                    if media_type in AV_MEDIA_TYPES:
                        totals["av_found"] += 1
                    elif media_type == "document":
                        totals["documents_found"] += 1

                if not media_files:
                    continue

                rows, skipped = await bulk_create_scraped_media(
                    db,
                    school=school,
                    source_page_url=scrape_url.url,
                    media_files=media_files,
                )
                totals["created"] += len(rows)
                totals["skipped"] += skipped

                # Enqueue ONLY newly created rows. This is what stops a
                # re-crawl from re-paying for the whole corpus, and it is safe
                # here because bulk_create_scraped_media has committed.
                for row in rows:
                    ingest_scraped_media.delay(row.id)
                    totals["enqueued"] += 1

    if failed_urls:
        logger.warning(
            "sweep_school_media: %s district(s) could not be scraped: %s",
            len(failed_urls),
            failed_urls,
        )
    logger.info("sweep_school_media finished: %s", totals)
    return totals


async def _scrape_one_batch_url(url: str, crawl_depth: int) -> dict:
    """Run one scrape attempt with no DB access, so it is safe to run
    concurrently alongside the other URLs in the same batch."""
    import httpx

    from app.services.web_scraper.school_scraper_service import SchoolScraperService

    async with SchoolScraperService() as scraper:
        try:
            result = await scraper.scrape_media_files(
                page_url=url, crawl_depth=crawl_depth
            )
        except httpx.HTTPStatusError as exc:
            return {
                "success": False,
                "http_status": exc.response.status_code if exc.response else None,
                "pages_crawled": 0,
                "media_files": [],
            }
        except Exception:
            # One bad URL must not sink the whole batch — every other URL in
            # the gather() below keeps running regardless.
            logger.exception("Batch scrape failed for %s", url)
            return {
                "success": False,
                "http_status": None,
                "pages_crawled": 0,
                "media_files": [],
            }

    return {
        "success": True,
        "http_status": 200,
        "pages_crawled": result["pages_crawled"],
        "media_files": result["media_files"],
    }


@celery_app.task(
    name="app.tasks.school_scraper_tasks.scrape_media_batch",
    bind=True,
    max_retries=3,
)
def scrape_media_batch(
    self,
    tenant_id: int,
    school_id: int,
    scrape_url_ids: list[int],
    crawl_depth: int,
) -> dict:
    """Background counterpart of POST /school-scraper/scrape-media-batch.

    Crawls every scrape_url_id concurrently (network I/O only, no DB), then
    persists sequentially — an AsyncSession cannot safely be touched from
    concurrent coroutines. Poll GET /scrape-media-batch/status?task_id=...
    for the result.

    Always persists discovered media, unlike /scrape-media's optional
    persist flag: this task backs only "Scrape selected" on the Schools
    admin page, which has no preview-only use case.
    """
    try:
        loop = get_event_loop()
        return loop.run_until_complete(
            _scrape_media_batch_async(
                tenant_id, school_id, scrape_url_ids, crawl_depth
            )
        )
    except Exception as exc:
        logger.error("scrape_media_batch failed: %s", str(exc))
        raise self.retry(
            exc=exc, countdown=60 * (2**self.request.retries)
        ) from exc


async def _scrape_media_batch_async(
    tenant_id: int,
    school_id: int,
    scrape_url_ids: list[int],
    crawl_depth: int,
) -> dict:
    import asyncio

    from app.crud.schools import (
        bulk_create_scraped_media,
        get_school,
        record_scrape_result,
    )

    async with AsyncSessionLocal() as db:
        school = await get_school(db, tenant_id, school_id)
        if not school:
            return {
                "error": f"School {school_id} not found in tenant {tenant_id}",
                "results": [],
            }

        urls_by_id = {u.id: u for u in school.scrape_urls}
        missing = [sid for sid in scrape_url_ids if sid not in urls_by_id]
        if missing:
            return {
                "error": f"Scrape URL(s) {missing} not found for school {school_id}",
                "results": [],
            }

        ordered_urls = [urls_by_id[sid] for sid in scrape_url_ids]

        # The only part worth parallelizing: real network crawling. It has
        # no DB dependency, so every URL runs concurrently instead of the
        # sequential loop the old synchronous endpoint used — wall time
        # drops from sum(each URL) to roughly max(each URL).
        raw_results = await asyncio.gather(
            *(_scrape_one_batch_url(u.url, crawl_depth) for u in ordered_urls)
        )

        results: list[dict] = []
        for scrape_url, raw in zip(ordered_urls, raw_results, strict=True):
            # All DB writes happen here, sequentially, after every crawl has
            # finished — AsyncSession is not safe for concurrent use.
            await record_scrape_result(
                db,
                scrape_url,
                http_status=raw["http_status"],
                page_count=raw["pages_crawled"] if raw["success"] else None,
            )

            entry = {
                "source_url": scrape_url.url,
                "scrape_url_id": scrape_url.id,
                "success": raw["success"],
                "http_status": raw["http_status"],
                "pages_crawled": raw["pages_crawled"],
                "media_files": raw["media_files"],
                "persisted": 0,
                "skipped_duplicates": 0,
                "enqueued": 0,
                "scraped_media_ids": [],
            }

            if raw["success"] and raw["media_files"]:
                rows, skipped = await bulk_create_scraped_media(
                    db,
                    school=school,
                    source_page_url=scrape_url.url,
                    media_files=raw["media_files"],
                )
                entry["persisted"] = len(rows)
                entry["skipped_duplicates"] = skipped
                entry["scraped_media_ids"] = [r.id for r in rows]
                # Enqueued only AFTER the commit inside bulk_create_scraped_media
                # — otherwise the worker's db.get(ScrapedMedia, id) returns None
                # and silently drops the item.
                for row in rows:
                    ingest_scraped_media.delay(row.id)
                    entry["enqueued"] += 1

            results.append(entry)

    return {"error": None, "results": results}
