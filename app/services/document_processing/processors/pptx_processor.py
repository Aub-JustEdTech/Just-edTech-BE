"""PPTX (PowerPoint) document processor"""

import logging
from typing import Any

try:
    from pptx import Presentation
except ImportError:
    Presentation = None  # type: ignore[assignment]

from app.services.document_processing.base import DocumentProcessor

logger = logging.getLogger(__name__)


class PPTXProcessor(DocumentProcessor):
    """Process PPTX (Microsoft PowerPoint) documents via python-pptx.

    Slides are exposed both as a flat text blob (via `extract_text`) and as
    a per-slide list (via `extract_text_by_page`) so the ingest pipeline can
    attach `page_number = slide_index` to chunks — matching how PDF chunks
    preserve their page number.
    """

    supported_extensions = [".pptx"]
    supported_mime_types = [
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ]

    def extract_text(self, file_path: str) -> str:
        """Extract concatenated text from all slides of a PPTX."""
        if Presentation is None:
            raise ImportError(
                "python-pptx is required for PPTX processing. "
                "Install it with: pip install python-pptx"
            )
        try:
            presentation = Presentation(file_path)
            text_parts: list[str] = []
            for slide in presentation.slides:
                slide_text = self._extract_slide_text(slide)
                if slide_text:
                    text_parts.append(slide_text)
            text = "\n\n".join(text_parts)
            logger.info(f"Extracted {len(text)} characters from PPTX")
            return text.strip()
        except Exception as e:
            logger.error(f"Error extracting text from PPTX: {e}")
            raise

    def extract_text_by_page(self, file_path: str) -> list[str]:
        """
        Extract text from a PPTX as a list of slide strings (1:1 with slides).

        Used to preserve `page_number` (slide index) metadata for downstream
        chunking and citations, mirroring PDFProcessor.extract_text_by_page.
        """
        if Presentation is None:
            raise ImportError(
                "python-pptx is required for PPTX processing. "
                "Install it with: pip install python-pptx"
            )
        try:
            presentation = Presentation(file_path)
            pages: list[str] = []
            for slide in presentation.slides:
                pages.append(self._extract_slide_text(slide).strip())
            return pages
        except Exception as e:
            logger.error(f"Error extracting per-slide text from PPTX: {e}")
            raise

    def extract_metadata(self, file_path: str) -> dict[str, Any]:
        """Extract PPTX metadata (slide count, core properties)."""
        if Presentation is None:
            raise ImportError(
                "python-pptx is required for PPTX processing. "
                "Install it with: pip install python-pptx"
            )
        try:
            presentation = Presentation(file_path)
            slide_count = len(presentation.slides)

            # python-pptx exposes core_properties similar to python-docx.
            core_props = presentation.core_properties
            metadata = {
                "slide_count": slide_count,
                "page_count": slide_count,  # alias for chunking code that
                # reads page_count generically.
                "author": getattr(core_props, "author", "") or "",
                "title": getattr(core_props, "title", "") or "",
                "subject": getattr(core_props, "subject", "") or "",
                "created": (
                    str(core_props.created) if getattr(core_props, "created", None) else ""
                ),
                "modified": (
                    str(core_props.modified) if getattr(core_props, "modified", None) else ""
                ),
            }
            return metadata
        except Exception as e:
            logger.error(f"Error extracting PPTX metadata: {e}")
            return {"slide_count": 0, "page_count": 0}

    def validate(self, file_path: str) -> bool:
        """Validate PPTX file by attempting to open it."""
        if Presentation is None:
            return False
        try:
            Presentation(file_path)
            return True
        except Exception:
            return False

    # ------------------------------------------------------------------
    # Internal helpers
    # ------------------------------------------------------------------

    @staticmethod
    def _extract_slide_text(slide: Any) -> str:
        """
        Extract every text-bearing run from a slide, including text frames
        inside shapes (autoshapes, placeholders) and tables.
        """
        parts: list[str] = []
        for shape in slide.shapes:
            # Text frame
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    run_text = "".join(run.text for run in paragraph.runs)
                    if not run_text:
                        # Fall back to paragraph-level text (some shapes
                        # only populate paragraph.text).
                        run_text = paragraph.text
                    if run_text.strip():
                        parts.append(run_text)
            # Table cells
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            parts.append(cell.text)
        return "\n".join(parts)
